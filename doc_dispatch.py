#!/usr/bin/env python3
"""
doc_dispatch.py — File-type dispatcher for dropped documents.
Usage: python3 doc_dispatch.py <file> [file2 ...]
       or drag-and-drop onto the script in Terminal
Outputs extracted text to <filename>.txt alongside the source file.
"""

import sys
import os

def extract_pdf(path):
    import fitz
    doc = fitz.open(path)
    return "\n\n".join(page.get_text() for page in doc)

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_odt(path):
    from odf.opendocument import load
    from odf.text import P
    doc = load(path)
    return "\n".join(
        str(p) for p in doc.text.getElementsByType(P) if str(p).strip()
    )

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)

def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"--- Sheet: {sheet} ---")
        for row in ws.iter_rows(values_only=True):
            row_str = "\t".join(str(c) if c is not None else "" for c in row)
            if row_str.strip():
                lines.append(row_str)
    return "\n".join(lines)

def extract_xls(path):
    import xlrd
    wb = xlrd.open_workbook(path)
    lines = []
    for sheet in wb.sheets():
        lines.append(f"--- Sheet: {sheet.name} ---")
        for row in range(sheet.nrows):
            row_str = "\t".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols))
            if row_str.strip():
                lines.append(row_str)
    return "\n".join(lines)

def extract_ods(path):
    from odf.opendocument import load
    from odf.table import Table, TableRow, TableCell
    from odf.text import P
    doc = load(path)
    lines = []
    for table in doc.spreadsheet.getElementsByType(Table):
        lines.append(f"--- Sheet: {table.getAttribute('name')} ---")
        for row in table.getElementsByType(TableRow):
            cells = []
            for cell in row.getElementsByType(TableCell):
                ps = cell.getElementsByType(P)
                cells.append(" ".join(str(p) for p in ps))
            line = "\t".join(cells)
            if line.strip():
                lines.append(line)
    return "\n".join(lines)

def extract_image(path):
    import pytesseract
    from PIL import Image
    img = Image.open(path)
    return pytesseract.image_to_string(img)

def extract_html(path):
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    return soup.get_text(separator="\n").strip()

def extract_text(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_csv(path):
    import pandas as pd
    df = pd.read_csv(path)
    return df.to_string(index=False)

DISPATCH = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".odt":  extract_odt,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xls,
    ".ods":  extract_ods,
    ".png":  extract_image,
    ".jpg":  extract_image,
    ".jpeg": extract_image,
    ".tiff": extract_image,
    ".tif":  extract_image,
    ".bmp":  extract_image,
    ".html": extract_html,
    ".htm":  extract_html,
    ".txt":  extract_text,
    ".md":   extract_text,
    ".csv":  extract_csv,
}

def process_file(path):
    ext = os.path.splitext(path)[1].lower()
    handler = DISPATCH.get(ext)
    if not handler:
        print(f"[SKIP] Unsupported type: {path}")
        return

    print(f"[→] Processing {os.path.basename(path)} ({ext})")
    try:
        text = handler(path)
        out_path = os.path.splitext(path)[0] + "_extracted.txt"
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"[✓] Saved → {out_path}")
    except Exception as e:
        print(f"[✗] Failed: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 doc_dispatch.py <file> [file2 ...]")
        sys.exit(1)
    for arg in sys.argv[1:]:
        if os.path.isfile(arg):
            process_file(arg)
        else:
            print(f"[SKIP] Not a file: {arg}")
